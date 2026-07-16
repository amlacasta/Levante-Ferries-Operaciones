from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "deliverables/Proyecto_Control_Operativo_Ferries.pptx"

BG = RGBColor(8, 17, 31)
PANEL = RGBColor(17, 32, 52)
WHITE = RGBColor(235, 243, 255)
MUTED = RGBColor(151, 170, 195)
CYAN = RGBColor(83, 223, 209)
AMBER = RGBColor(255, 200, 87)

KPI = pd.read_csv(ROOT / "reports/tables/executive_kpis.csv").iloc[0]
ROUTES = pd.read_csv(ROOT / "reports/tables/route_scorecard.csv")
CAUSES = pd.read_csv(ROOT / "reports/tables/cause_pareto.csv")
VESSELS = pd.read_csv(ROOT / "reports/tables/vessel_rotation_scorecard.csv")
SCENARIOS = pd.read_csv(ROOT / "reports/tables/scenarios.csv")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_text(slide, text, x, y, w, h, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.alignment = align
    run = p.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def base_slide(title, kicker):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_text(slide, kicker.upper(), .6, .25, 8, .3, 10, CYAN, True)
    add_text(slide, title, .6, .58, 12, .55, 27, WHITE, True)
    return slide


def bullets(slide, items, x=.75, y=1.55, w=11.8, h=4.8, size=18):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)


def table(slide, df, x=.7, y=1.5, w=12, h=4.9, size=10):
    rows, cols = len(df) + 1, len(df.columns)
    t = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for c, name in enumerate(df.columns):
        t.cell(0, c).text = str(name)
        t.cell(0, c).fill.solid(); t.cell(0, c).fill.fore_color.rgb = RGBColor(42, 62, 88)
    for r, row in enumerate(df.itertuples(index=False), 1):
        for c, value in enumerate(row):
            t.cell(r, c).text = str(value)
            t.cell(r, c).fill.solid(); t.cell(r, c).fill.fore_color.rgb = PANEL
    for row in t.rows:
        for cell in row.cells:
            for p in cell.text_frame.paragraphs:
                p.font.name = "Aptos"; p.font.size = Pt(size); p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER


slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid(); slide.background.fill.fore_color.rgb = BG
add_text(slide, "OPERATIONS ANALYTICS · PORTFOLIO", .7, .7, 8, .35, 11, CYAN, True)
add_text(slide, "Levante-Ferries:\nOperaciones", .7, 1.25, 7, 1.4, 38, WHITE, True)
add_text(slide, "Flota física · rotaciones · KPIs · diagnóstico · escenarios · simulador", .7, 3.0, 8, .7, 18, MUTED)
add_text(slide, f"{int(KPI.scheduled_services):,}".replace(',', '.'), 9.2, 1.35, 3, .6, 34, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, "servicios sintéticos", 9.2, 1.95, 3, .3, 13, MUTED, False, PP_ALIGN.CENTER)
add_text(slide, "CASO FICTICIO · DATOS SINTÉTICOS", .7, 6.7, 6, .3, 11, AMBER, True)

slide = base_slide("Resumen ejecutivo", "01 · Executive summary")
bullets(slide, [
    f"12.000 servicios programados y {int(KPI.completed_services):,} completados".replace(',', '.'),
    f"Fiabilidad operacional: {KPI.operational_reliability_pct:.2f}%",
    f"OTR≤15 de completados: {KPI.otr_15_completed_pct:.2f}%",
    f"Cancel Rate: {KPI.cancel_rate_pct:.2f}%",
    f"Retraso propagado: {KPI.total_propagated_delay_min:,.0f} min".replace(',', '.'),
    f"Margen total simulado: {KPI.total_margin_eur/1_000_000:.1f} M€",
])

slide = base_slide("Fiabilidad por ruta", "02 · Route scorecard")
r = ROUTES.sort_values("operational_reliability_pct").head(8).copy()
r = r[["route_id","scheduled_services","cancel_rate_pct","otr_15_completed_pct","operational_reliability_pct","avg_delay_min","status"]]
r.columns = ["Ruta","Servicios","Cancel %","OTR15 %","Fiabilidad %","Delay","Estado"]
table(slide, r.round(2))

slide = base_slide("Principales causas", "03 · Delay Pareto")
c = CAUSES.head(6)[["disruption_reason","services","total_delay_min","delay_share_pct","p95_delay_min","controllability_score"]].copy()
c.columns = ["Causa","Servicios","Delay total","Share %","P95","Control"]
table(slide, c.round(2))

slide = base_slide("Rotaciones y activos", "04 · Fleet propagation")
v = VESSELS.head(10)[["vessel_id","vessel_type","scheduled_services","avg_delay_min","rotation_impacted_pct","total_propagated_delay_min"]].copy()
v.columns = ["Buque","Tipo","Servicios","Delay","Impactado %","Delay propagado"]
table(slide, v.round(2))

slide = base_slide("Escenarios de sensibilidad", "05 · What-if")
s = SCENARIOS[["scenario","delay_saved_min","delay_saved_pct","direct_delay_cost_saved_eur"]].copy()
s.columns = ["Escenario","Min ahorrados","Ahorro %","Coste evitado €"]
table(slide, s.round(2), size=11)

slide = base_slide("Recomendaciones", "06 · Action plan")
bullets(slide, [
    "Protocolos dinámicos para WINDY, ROUGH y STORM.",
    "Mantenimiento preventivo en buques con mayor propagación.",
    "Revisión de tiempos de escala y capacidad de recuperación.",
    "Coordinación portuaria en rutas y horas críticas.",
    "Control Tower semanal con fiabilidad, costes y hotspots.",
])

slide = base_slide("Limitaciones y siguientes pasos", "07 · Roadmap")
bullets(slide, [
    "Datos, flota, demanda, meteorología y costes completamente sintéticos.",
    "Los escenarios son sensibilidades, no estimaciones causales.",
    "Siguiente evolución: datos reales, forecasting, pricing y optimización de rotaciones.",
    "Integración futura con Power BI, Streamlit y alertas operativas.",
])

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(OUT)
